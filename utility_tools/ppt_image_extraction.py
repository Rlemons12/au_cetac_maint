from pptx import Presentation
import os

def extract_images_from_ppt(ppt_file, output_dir):
    # Load the PowerPoint presentation
    prs = Presentation(ppt_file)

    image_count = 0

    # Iterate through each slide in the presentation
    for slide in prs.slides:
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape is an image
            if shape.shape_type == 13:  # 13 corresponds to a picture shape
                image_count += 1
                # Extract the image
                image = shape.image
                # Save the image to a file
                image_bytes = image.blob
                with open(os.path.join(output_dir, f"image_{image_count}.png"), "wb") as f:
                    f.write(image_bytes)
    
    print(f"Extracted {image_count} images from {ppt_file} to {output_dir}")

# Provide the path to your PowerPoint file
ppt_file = r"C:\Users\10169062\Desktop\AI_EMTACr7c6\Docs\TTP_CPU_CARD_OPERATION_rev_B.pptx"

# Provide the output directory
output_dir = r"C:\\Users\\10169062\\Desktop\\AI_EMTACr7c6\\Docs\\extracted images"

# Call the function to extract images
extract_images_from_ppt(ppt_file, output_dir)
