import os
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image

# Function to convert PPT to PDF with images
def convert_ppt_to_pdf(ppt_file):
    prs = Presentation(ppt_file)
    pdf_file = os.path.splitext(ppt_file)[0] + ".pdf"  # Use the same filename with .pdf extension

    pdf = canvas.Canvas(pdf_file, pagesize=letter)

    for slide in prs.slides:
        pdf.showPage()
        pdf.setFont("Helvetica", 12)
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.drawString(100, 500, shape.text)
                
            if shape.has_image:
                img = shape.image
                img_bytes = img.blob
                img_file = "temp_img.png"
                
                with open(img_file, "wb") as f:
                    f.write(img_bytes)
                
                img = Image.open(img_file)
                pdf.drawImage(img_file, 100, 300, width=400, height=300)
                os.remove(img_file)

    pdf.save()

# Folder containing PPT files
ppt_folder = r'C:\Users\10169062\Desktop\AI_EMTACr3r4\PPT_FILES'

# Convert each PPT file in the folder to PDF
for filename in os.listdir(ppt_folder):
    if filename.endswith(".pptx"):
        ppt_file = os.path.join(ppt_folder, filename)
        convert_ppt_to_pdf(ppt_file)
        print(f"Converted {filename} to PDF")

print("Conversion complete! PDF files with images are saved in the same folder.")
